import os
from typing import List, Optional
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import InMemoryVectorStore
from langchain_ollama import OllamaEmbeddings
from langchain_core.documents import Document


class KnowledgeBase:
    def __init__(
        self,
        embedding_model: str = "nomic-embed-text",
        base_url: str = "http://localhost:11434",
        chunk_size: int = 500,
        chunk_overlap: int = 100,
    ):
        self.embedding_model = embedding_model
        self.base_url = base_url
        self.embeddings = OllamaEmbeddings(
            model=embedding_model,
            base_url=base_url,
        )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
        self.vectorstore = None

    def _load_pptx(self, file_path: str) -> List[Document]:
        from pptx import Presentation

        prs = Presentation(file_path)
        full_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if parts:
                full_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

        if not full_text:
            return []
        return [
            Document(
                page_content="\n\n".join(full_text),
                metadata={"source": os.path.basename(file_path)},
            )
        ]

    def load_document(self, file_path: str) -> List[Document]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
        elif ext == ".docx":
            loader = Docx2txtLoader(file_path)
        elif ext in [".txt", ".md", ".markdown"]:
            loader = TextLoader(file_path, encoding="utf-8")
        elif ext == ".pptx":
            return self._load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

        return loader.load()

    def add_documents(self, file_paths: List[str]) -> int:
        all_documents = []
        for file_path in file_paths:
            try:
                documents = self.load_document(file_path)
                for doc in documents:
                    doc.metadata["source"] = os.path.basename(file_path)
                all_documents.extend(documents)
            except Exception as e:
                print(f"[KB] 加载文件 {file_path} 失败: {e}")
                continue

        if not all_documents:
            return 0

        split_docs = self.text_splitter.split_documents(all_documents)

        if self.vectorstore is None:
            self.vectorstore = InMemoryVectorStore.from_documents(
                documents=split_docs,
                embedding=self.embeddings,
            )
        else:
            self.vectorstore.add_documents(split_docs)

        return len(split_docs)

    def search(self, query: str, k: int = 3) -> List[Document]:
        if self.vectorstore is None:
            return []
        return self.vectorstore.similarity_search(query, k=k)

    def get_retriever(self, k: int = 3):
        if self.vectorstore is None:
            return None
        return self.vectorstore.as_retriever(search_kwargs={"k": k})

    def get_chunks_count(self) -> int:
        if self.vectorstore is None:
            return 0
        try:
            return len(self.vectorstore.store)
        except Exception:
            return 0

    def clear_database(self):
        self.vectorstore = None

    def get_sources(self) -> List[str]:
        if self.vectorstore is None:
            return []
        try:
            if hasattr(self.vectorstore, 'store'):
                docs = self.vectorstore.store
                sources = set()
                for doc in docs.values():
                    if isinstance(doc, dict):
                        meta = doc.get('metadata') or {}
                        if isinstance(meta, dict) and 'source' in meta:
                            sources.add(meta['source'])
                    elif hasattr(doc, 'metadata'):
                        meta = doc.metadata
                        if isinstance(meta, dict) and 'source' in meta:
                            sources.add(meta['source'])
                return list(sources)
            return []
        except Exception as e:
            print(f"[KB] get_sources error: {e}")
            return []


def main():
    kb = KnowledgeBase()
    print(f"chunks: {kb.get_chunks_count()}, sources: {kb.get_sources()}")


if __name__ == "__main__":
    main()
